import pytesseract
from pdf2image import convert_from_path
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
from PIL import Image
import tkinter as tk
from tkinter import filedialog
import os

# set tesseract path
pytesseract.pytesseract.tesseract_cmd = r"D:\tessaract\tesseract.exe"

output_file = "output.txt"


def read_pdf(path):

    print("Reading PDF...")

    text = extract_text(path)

    if text and len(text.strip()) > 50:
        print("Text PDF detected")
        return text

    print("Scanned PDF detected → running OCR")

    images = convert_from_path(path)

    text = ""

    for i, img in enumerate(images):
        print(f"OCR page {i+1}/{len(images)}")
        text += pytesseract.image_to_string(img)

    return text


def read_docx(path):

    print("Reading DOCX...")

    doc = Document(path)

    return "\n".join(p.text for p in doc.paragraphs)


def read_pptx(path):

    print("Reading PPTX...")

    prs = Presentation(path)

    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


def read_image(path):

    print("Reading Image with OCR...")

    img = Image.open(path)

    return pytesseract.image_to_string(img)


def read_file(path):

    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        return read_pdf(path)

    elif ext in [".doc", ".docx"]:
        return read_docx(path)

    elif ext in [".ppt", ".pptx"]:
        return read_pptx(path)

    elif ext in [".png", ".jpg", ".jpeg"]:
        return read_image(path)

    else:
        raise ValueError("Unsupported file type")


def pick_and_read_files():
    root = tk.Tk()
    root.withdraw()
    
    # Needs to stay on top because some frameworks misbehave if UI is completely withdrawn without update
    root.update()
    root.wm_attributes('-topmost', 1)

    file_paths = filedialog.askopenfilenames(
        title="Select documents",
        filetypes=[
            ("Documents", "*.pdf *.doc *.docx *.ppt *.pptx *.png *.jpg *.jpeg")
        ]
    )

    results = []
    if file_paths:
        for file_path in file_paths:
            print("Selected file:", file_path)
            try:
                text = read_file(file_path)
                
                # Delete PDF if it was a PDF as requested
                if file_path.lower().endswith(".pdf"):
                    try:
                        os.remove(file_path)
                        print(f"Deleted original PDF file: {file_path}")
                    except Exception as e:
                        print(f"Could not delete PDF file: {e}")
                        
                if file_path and text:
                    results.append((file_path, text))
            except Exception as e:
                print(f"Error reading file {file_path}: {e}")
        return results
    else:
        print("No files selected")
        return []

if __name__ == "__main__":
    results = pick_and_read_files()
    if results:
        count = 0
        for path, text in results:
            with open(f"output_{count}.txt", "w", encoding="utf-8") as f:
                f.write(text)
            count += 1
        print(f"Successfully saved {count} files.")