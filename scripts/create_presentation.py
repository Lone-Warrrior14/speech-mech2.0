from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "SPEECH-MECH NEURAL ECOSYSTEM v2.0"
    subtitle.text = "Restructured Unified AI Intelligence\nPresented by: Lone-Warrior14"

    # Slide 2: Project Overview
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "PROJECT OVERVIEW"
    tf = slide.placeholders[1].text_frame
    tf.text = "A powerful, unified AI ecosystem featuring:"
    p = tf.add_paragraph()
    p.text = "• Biometric Secure Access"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• RAG Intelligence Memory (Pinecone)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Autonomous Image Synthesis (NVIDIA/Cloudflare)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Natural Language Voice UI"
    p.level = 1

    # Slide 3: Task-Based Architecture
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "RESTRUCTURED ARCHITECTURE v2.0"
    tf = slide.placeholders[1].text_frame
    tf.text = "Migrated from monolithic to task-based modules:"
    p = tf.add_paragraph()
    p.text = "• /authorization/: Identity & Biometrics"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• /rag_system/: Long-term Intelligence"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• /image gen/: Visual AI Backend"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• /speech_assistant/: Core Conversational Brain"
    p.level = 1

    # Slide 4: Biometric Neural Unlock
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "BIOMETRIC NEURAL UNLOCK"
    tf = slide.placeholders[1].text_frame
    tf.text = "State-of-the-art security layer:"
    p = tf.add_paragraph()
    p.text = "• Automated Neural Capturing for new users."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Real-time Anti-Spoofing (Liveness Check)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Multi-modal Identity Management (Face + Password)."
    p.level = 1

    # Slide 5: RAG Intelligence
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "RAG INTELLIGENCE MEMORY"
    tf = slide.placeholders[1].text_frame
    tf.text = "Augmenting LLMs with document knowledge:"
    p = tf.add_paragraph()
    p.text = "• Universal Reader (PDF, Audio, Docx)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Vectorized retrieval via Pinecone."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Context-aware responses for personal workspaces."
    p.level = 1

    # Slide 6: Continuous Visual Synthesis
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "CONTINUOUS VISUAL SYNTHESIS"
    tf = slide.placeholders[1].text_frame
    tf.text = "Advanced image generation workflow:"
    p = tf.add_paragraph()
    p.text = "• Primary: NVIDIA Flux.2 (High-Fidelity)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Backup: Cloudflare Workers (Fast-Uplink)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Direct API Integration into the Dashboard."
    p.level = 1

    # Slide 7: The Unified Launch System
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "UNIFIED HUB & DEPLOYMENT"
    tf = slide.placeholders[1].text_frame
    tf.text = "Seamless User Experience:"
    p = tf.add_paragraph()
    p.text = "• Unified Launcher (run.py) handles all components."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Secure Tunneling (ngrok) for mobile/phone access."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Professional Version Control (GitHub v2.0)."
    p.level = 1

    # Slide 8: Thank You
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "THANK YOU"
    p = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2)).text_frame.add_paragraph()
    p.text = "SPEECH-MECH v2.0 - The Future of Unified AI"
    p.font.size = Pt(36)
    p.alignment = PP_ALIGN.CENTER

    prs.save('SPEECH_MECH_Presentation.pptx')
    print("Presentation created successfully: SPEECH_MECH_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
