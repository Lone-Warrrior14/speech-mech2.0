import pytesseract
from pdf2image import convert_from_path
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
from PIL import Image
import tkinter as tk
from tkinter import filedialog
import os

# set tesseract path
pytesseract.pytesseract.tesseract_cmd = r"D:\tessaract\tesseract.exe"

def read_pdf(path):
    print("Reading PDF...")
    text = extract_text(path)
    if text and len(text.strip()) > 50:
        print("Text PDF detected")
        return text
    print("Scanned PDF detected → running OCR")
    try:
        images = convert_from_path(path)
        text = ""
        for i, img in enumerate(images):
            print(f"OCR page {i+1}/{len(images)}")
            text += pytesseract.image_to_string(img)
        return text
    except Exception as e:
        print(f"OCR Error: {e}")
        return ""

def read_docx(path):
    print("Reading DOCX...")
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)

def read_pptx(path):
    print("Reading PPTX...")
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_image(path):
    print("Reading Image with OCR...")
    img = Image.open(path)
    return pytesseract.image_to_string(img)

def read_audio(path):
    import speech_recognition as sr
    print(f"Transcribing audio: {path}...")
    r = sr.Recognizer()
    try:
        with sr.AudioFile(path) as source:
            audio = r.record(source)
        text = r.recognize_google(audio)
        print("Transcription complete.")
        return text
    except Exception as e:
        print(f"Transcription failed: {e}")
        return ""

def read_file(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf": return read_pdf(path)
    elif ext in [".doc", ".docx"]: return read_docx(path)
    elif ext in [".ppt", ".pptx"]: return read_pptx(path)
    elif ext in [".png", ".jpg", ".jpeg"]: return read_image(path)
    elif ext in [".wav", ".flac", ".mp3"]: return read_audio(path)
    else: raise ValueError("Unsupported file type")

def pick_and_read_files(media_only=False):
    root = tk.Tk()
    root.withdraw()
    root.update(); root.lift(); root.attributes('-topmost', True); root.focus_force()

    if media_only:
        file_paths = filedialog.askopenfilenames(
            title="Select Media for Intel Extraction",
            filetypes=[("Audio/Video", "*.wav *.flac *.mp3 *.mp4 *.wmv")]
        )
    else:
        file_paths = filedialog.askopenfilenames(
            title="Select Documents",
            filetypes=[("Docs", "*.pdf *.docx *.pptx *.png *.jpg *.jpeg *.wav *.mp3")]
        )

    results = []
    if file_paths:
        for file_path in file_paths:
            try:
                text = read_file(file_path)
                if text: results.append((file_path, text))
            except Exception as e:
                print(f"Error processing {file_path}: {e}")
        root.destroy()
        return results
    root.destroy()
    return []

if __name__ == "__main__":
    pass